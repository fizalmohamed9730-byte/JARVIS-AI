"""Document generator for creating Word, Excel, PowerPoint, PDF, and CSV files."""

import csv
import io
import logging
import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Union

logger = logging.getLogger(__name__)


class DocumentGenerator:
    """Generates office documents and common file formats programmatically.

    Uses openpyxl for Excel, python-docx for Word, python-pptx for
    presentations, and reportlab for PDF generation. Falls back to
    simpler implementations when optional dependencies are unavailable.
    """

    # --------------------------------------------------------------------- #
    # Word (.docx)
    # --------------------------------------------------------------------- #

    def generate_word(
        self,
        content: str,
        title: str = "",
        output_path: str = "",
    ) -> Dict[str, Any]:
        """Generate a Word document.

        Args:
            content: The document body text. Paragraphs separated by newlines.
            title: Optional title heading at the top of the document.
            output_path: Where to save the file. Auto-generated if empty.
        """
        output_path = output_path or self._default_path("document", ".docx")
        self._ensure_dir(output_path)

        try:
            from docx import Document
            from docx.shared import Pt, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()

            if title:
                heading = doc.add_heading(title, level=0)
                heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
                doc.add_paragraph("")

            for paragraph in content.split("\n\n"):
                stripped = paragraph.strip()
                if not stripped:
                    continue
                if stripped.startswith("# "):
                    doc.add_heading(stripped[2:], level=1)
                elif stripped.startswith("## "):
                    doc.add_heading(stripped[3:], level=2)
                elif stripped.startswith("### "):
                    doc.add_heading(stripped[4:], level=3)
                elif stripped.startswith("- "):
                    for item in stripped.split("\n"):
                        item = item.strip()
                        if item.startswith("- "):
                            doc.add_paragraph(item[2:], style="List Bullet")
                else:
                    doc.add_paragraph(stripped)

            doc.save(output_path)
            logger.info("Generated Word document: %s", output_path)
            return {"success": True, "path": output_path}
        except ImportError:
            # Fallback: save as plain text
            return self._fallback_text(content, output_path, ".docx")
        except Exception as exc:
            logger.exception("Failed to generate Word document")
            return {"success": False, "error": str(exc)}

    # --------------------------------------------------------------------- #
    # Excel (.xlsx)
    # --------------------------------------------------------------------- #

    def generate_excel(
        self,
        data: Sequence[Sequence[Any]],
        headers: Optional[Sequence[str]] = None,
        output_path: str = "",
        sheet_name: str = "Sheet1",
    ) -> Dict[str, Any]:
        """Generate an Excel spreadsheet.

        Args:
            data: 2D list of row data.
            headers: Optional header row.
            output_path: Where to save.
            sheet_name: Name of the worksheet.
        """
        output_path = output_path or self._default_path("spreadsheet", ".xlsx")
        self._ensure_dir(output_path)

        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment

            wb = Workbook()
            ws = wb.active
            ws.title = sheet_name

            # Write headers
            if headers:
                header_font = Font(bold=True, color="FFFFFF")
                header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                for col_idx, header in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=col_idx, value=header)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = Alignment(horizontal="center")

            # Write data rows
            start_row = 2 if headers else 1
            for row_idx, row_data in enumerate(data, start_row):
                for col_idx, value in enumerate(row_data, 1):
                    ws.cell(row=row_idx, column=col_idx, value=value)

            # Auto-fit column widths
            for col_idx, column in enumerate(ws.columns, 1):
                max_length = 0
                for cell in column:
                    try:
                        if cell.value:
                            max_length = max(max_length, len(str(cell.value)))
                    except Exception:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column[0].column_letter].width = adjusted_width

            wb.save(output_path)
            logger.info("Generated Excel file: %s", output_path)
            return {"success": True, "path": output_path, "rows": len(data)}
        except ImportError:
            # Fallback: save as CSV
            csv_path = output_path.replace(".xlsx", ".csv")
            return self._generate_csv_fallback(data, headers, csv_path)
        except Exception as exc:
            logger.exception("Failed to generate Excel file")
            return {"success": False, "error": str(exc)}

    # --------------------------------------------------------------------- #
    # PowerPoint (.pptx)
    # --------------------------------------------------------------------- #

    def generate_presentation(
        self,
        slides: Sequence[Dict[str, Any]],
        output_path: str = "",
    ) -> Dict[str, Any]:
        """Generate a PowerPoint presentation.

        Args:
            slides: List of slide dicts with keys:
                - ``title``: Slide title (required).
                - ``content``: Body text (optional).
                - ``notes``: Speaker notes (optional).
            output_path: Where to save.
        """
        output_path = output_path or self._default_path("presentation", ".pptx")
        self._ensure_dir(output_path)

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            title_layout = prs.slide_layouts[0]
            content_layout = prs.slide_layouts[1]

            for slide_data in slides:
                if slide_data.get("blank"):
                    blank_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_layout)
                else:
                    slide = prs.slides.add_slide(title_layout)
                    if slide.shapes.title:
                        slide.shapes.title.text = slide_data.get("title", "")
                    if len(slide.placeholders) > 1:
                        body = slide.placeholders[1]
                        body.text = slide_data.get("content", "")

                # Add content text box if there's no content layout
                if slide_data.get("content") and not slide_data.get("blank"):
                    pass  # Already handled by content layout

                # Add speaker notes
                if slide_data.get("notes"):
                    notes_slide = slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    text_frame.text = slide_data["notes"]

            prs.save(output_path)
            logger.info("Generated presentation: %s with %d slides", output_path, len(slides))
            return {"success": True, "path": output_path, "slide_count": len(slides)}
        except ImportError:
            return {"success": False, "error": "python-pptx is required: pip install python-pptx"}
        except Exception as exc:
            logger.exception("Failed to generate presentation")
            return {"success": False, "error": str(exc)}

    # --------------------------------------------------------------------- #
    # PDF
    # --------------------------------------------------------------------- #

    def generate_pdf(
        self,
        content: str,
        output_path: str = "",
        title: str = "",
    ) -> Dict[str, Any]:
        """Generate a PDF document.

        Args:
            content: Text content. Supports basic markdown-like headings (# ## ###)
                     and lists (- item).
            output_path: Where to save.
            title: Optional document title.
        """
        output_path = output_path or self._default_path("document", ".pdf")
        self._ensure_dir(output_path)

        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem
            from reportlab.lib.enums import TA_CENTER

            doc = SimpleDocTemplate(
                output_path,
                pagesize=A4,
                rightMargin=inch,
                leftMargin=inch,
                topMargin=inch,
                bottomMargin=inch,
            )

            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                "CustomTitle",
                parent=styles["Title"],
                fontSize=24,
                alignment=TA_CENTER,
                spaceAfter=20,
            )
            heading1_style = ParagraphStyle(
                "Heading1Custom",
                parent=styles["Heading1"],
                fontSize=18,
                spaceBefore=16,
                spaceAfter=8,
            )
            heading2_style = ParagraphStyle(
                "Heading2Custom",
                parent=styles["Heading2"],
                fontSize=14,
                spaceBefore=12,
                spaceAfter=6,
            )
            body_style = styles["BodyText"]

            story = []

            if title:
                story.append(Paragraph(title, title_style))
                story.append(Spacer(1, 12))

            for paragraph in content.split("\n\n"):
                stripped = paragraph.strip()
                if not stripped:
                    continue

                if stripped.startswith("# "):
                    story.append(Paragraph(stripped[2:], heading1_style))
                elif stripped.startswith("## "):
                    story.append(Paragraph(stripped[3:], heading2_style))
                elif stripped.startswith("### "):
                    story.append(Paragraph(stripped[4:], heading2_style))
                elif stripped.startswith("- "):
                    items = [
                        ListItem(Paragraph(item.strip()[2:], body_style))
                        for item in stripped.split("\n")
                        if item.strip().startswith("- ")
                    ]
                    if items:
                        story.append(ListFlowable(items, bulletType="bullet"))
                else:
                    safe_text = stripped.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    story.append(Paragraph(safe_text, body_style))

                story.append(Spacer(1, 6))

            doc.build(story)
            logger.info("Generated PDF: %s", output_path)
            return {"success": True, "path": output_path}
        except ImportError:
            # Fallback: save as plain text
            txt_path = output_path.replace(".pdf", ".txt")
            return self._fallback_text(content, txt_path, ".pdf")
        except Exception as exc:
            logger.exception("Failed to generate PDF")
            return {"success": False, "error": str(exc)}

    # --------------------------------------------------------------------- #
    # CSV
    # --------------------------------------------------------------------- #

    def generate_csv(
        self,
        data: Sequence[Sequence[Any]],
        headers: Optional[Sequence[str]] = None,
        output_path: str = "",
    ) -> Dict[str, Any]:
        """Generate a CSV file.

        Args:
            data: 2D list of row data.
            headers: Optional header row.
            output_path: Where to save.
        """
        output_path = output_path or self._default_path("data", ".csv")
        self._ensure_dir(output_path)

        try:
            with open(output_path, "w", newline="", encoding="utf-8") as f:
                writer = csv.writer(f, quoting=csv.QUOTE_ALL)
                if headers:
                    writer.writerow(headers)
                for row in data:
                    writer.writerow(row)

            logger.info("Generated CSV: %s (%d rows)", output_path, len(data))
            return {"success": True, "path": output_path, "rows": len(data)}
        except Exception as exc:
            logger.exception("Failed to generate CSV")
            return {"success": False, "error": str(exc)}

    # --------------------------------------------------------------------- #
    # Helpers
    # --------------------------------------------------------------------- #

    def _fallback_text(
        self, content: str, output_path: str, original_ext: str
    ) -> Dict[str, Any]:
        """Save content as a text file when the target format library is unavailable."""
        try:
            txt_path = output_path.replace(original_ext, ".txt")
            with open(txt_path, "w", encoding="utf-8") as f:
                f.write(content)
            return {
                "success": True,
                "path": txt_path,
                "note": f"Saved as .txt (install library for {original_ext} support)",
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def _generate_csv_fallback(
        self,
        data: Sequence[Sequence[Any]],
        headers: Optional[Sequence[str]],
        output_path: str,
    ) -> Dict[str, Any]:
        """CSV fallback when openpyxl is not installed."""
        return self.generate_csv(data, headers, output_path)

    @staticmethod
    def _default_path(prefix: str, extension: str) -> str:
        """Generate a default output path in the user's Documents folder."""
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        documents = os.path.expanduser("~/Documents")
        os.makedirs(documents, exist_ok=True)
        return os.path.join(documents, f"{prefix}_{timestamp}{extension}")

    @staticmethod
    def _ensure_dir(path: str) -> None:
        """Create parent directories if they don't exist."""
        parent = os.path.dirname(path)
        if parent:
            os.makedirs(parent, exist_ok=True)
